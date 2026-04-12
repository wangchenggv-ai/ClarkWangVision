from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xAB)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=18, color=DARK_GRAY):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(8)
    return txBox

# ==================== P1: 封面 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)
add_textbox(slide, 1.5, 2.0, 10, 1.2, '近视防控的全生命周期管理', 40, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 3.2, 10, 0.8, '从临床验证到产品创新', 28, RGBColor(0xA0,0xC4,0xE8), False, PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 5.0, 10, 0.5, '王成  |  高视星设计总监', 20, RGBColor(0xCC,0xCC,0xCC), False, PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 5.6, 10, 0.5, '2026年4月9日  ·  眼视光+角膜塑形学术会议', 16, RGBColor(0x99,0x99,0x99), False, PP_ALIGN.CENTER)

# ==================== P2: 开场故事 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, 1.5, 1.0, 10, 0.6, 'P2', 14, MID_GRAY, False, PP_ALIGN.RIGHT)
add_textbox(slide, 1.5, 2.0, 10, 1.0, '一个看不清的字', 36, DARK_BLUE, True, PP_ALIGN.LEFT)

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(10), Inches(3.0))
tf = txBox.text_frame
tf.word_wrap = True
lines = [
    'OK镜包装上的溯源码，字非常小',
    '做研发第一年能看清，第二年看不清了',
    '',
    '近视 + 老花',
    '现有OK镜的设计参数里，没有一个是为这个问题准备的',
]
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(22)
    p.font.color.rgb = DARK_GRAY
    p.font.name = 'Microsoft YaHei'
    p.space_after = Pt(6)

# ==================== P3: 临床研究概况 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, 1.5, 1.0, 10, 0.6, 'P3', 14, MID_GRAY, False, PP_ALIGN.RIGHT)
add_textbox(slide, 1.5, 1.5, 10, 0.8, '临床研究概况', 32, DARK_BLUE, True)

# Data boxes
box_data = [
    ('211 例', '6-14岁近视儿童'),
    ('4 组', '平行对照'),
    ('6 个月', '随访周期'),
    ('豪雅新乐学', '对照组'),
]
for i, (big, small) in enumerate(box_data):
    left = 1.5 + i * 2.8
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(3.0), Inches(2.5), Inches(2.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = big
    p.font.size = Pt(32)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p2 = tf.add_paragraph()
    p2.text = small
    p2.font.size = Pt(16)
    p2.font.color.rgb = MID_GRAY
    p2.font.name = 'Microsoft YaHei'
    p2.alignment = PP_ALIGN.CENTER

add_textbox(slide, 1.5, 5.5, 10, 0.5, '石家庄人民医院 前瞻性对照研究', 14, MID_GRAY, False, PP_ALIGN.LEFT)

# ==================== P4: 核心数据 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, 1.5, 1.0, 10, 0.6, 'P4', 14, MID_GRAY, False, PP_ALIGN.RIGHT)
add_textbox(slide, 1.5, 1.5, 10, 0.8, '核心数据  ·  Ultra组', 32, DARK_BLUE, True)

# Big numbers
add_textbox(slide, 1.5, 2.8, 5, 1.5, '0.138', 96, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 4.2, 5, 0.6, 'mm/年 眼轴增长', 20, MID_GRAY, False, PP_ALIGN.CENTER)

add_textbox(slide, 7.0, 2.8, 5, 1.5, '72.5%', 96, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_textbox(slide, 7.0, 4.2, 5, 0.6, '相对减缓率', 20, MID_GRAY, False, PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 5.5, 10, 0.5, '国际同类研究第一梯队水平', 18, DARK_GRAY, True, PP_ALIGN.CENTER)

# ==================== P5: 低龄组发现 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, 1.5, 1.0, 10, 0.6, 'P5', 14, MID_GRAY, False, PP_ALIGN.RIGHT)
add_textbox(slide, 1.5, 1.5, 10, 0.8, '低龄组更需要早期干预', 32, DARK_BLUE, True)

# Comparison
add_textbox(slide, 2.0, 3.0, 4, 1.0, '7-9岁组', 24, DARK_GRAY, True, PP_ALIGN.CENTER)
add_textbox(slide, 2.0, 3.8, 4, 1.0, '3×', 72, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_textbox(slide, 2.0, 5.0, 4, 0.5, '眼轴增长速度', 16, MID_GRAY, False, PP_ALIGN.CENTER)

add_textbox(slide, 7.0, 3.5, 4, 0.5, 'vs', 36, MID_GRAY, False, PP_ALIGN.CENTER)

add_textbox(slide, 7.0, 3.0, 4, 1.0, '10-14岁组', 24, DARK_GRAY, True, PP_ALIGN.CENTER)
add_textbox(slide, 7.0, 3.8, 4, 1.0, '1×', 72, MID_GRAY, True, PP_ALIGN.CENTER)
add_textbox(slide, 7.0, 5.0, 4, 0.5, '眼轴增长速度', 16, MID_GRAY, False, PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 5.8, 10, 0.5, 'p = 0.0003  |  越早介入，效果越可量化', 18, DARK_GRAY, True, PP_ALIGN.CENTER)

# ==================== P6: 产品矩阵 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, 1.5, 1.0, 10, 0.6, 'P6', 14, MID_GRAY, False, PP_ALIGN.RIGHT)
add_textbox(slide, 1.5, 1.5, 10, 0.8, '产品矩阵扩展', 32, DARK_BLUE, True)

products = [
    ('远视储备', '正度数方案', '临床已入组'),
    ('成人OK镜', '38-55岁 近视+早老花', '注册证6月获批 · 7月上市'),
    ('离焦RGP &\n超级非球面离焦镜', '离焦+点扩散二合一', '行业首次融合设计'),
]
for i, (title, sub1, sub2) in enumerate(products):
    left = 1.5 + i * 3.7
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.8), Inches(3.4), Inches(3.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p2 = tf.add_paragraph()
    p2.text = sub1
    p2.font.size = Pt(16)
    p2.font.color.rgb = DARK_GRAY
    p2.font.name = 'Microsoft YaHei'
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    p3 = tf.add_paragraph()
    p3.text = sub2
    p3.font.size = Pt(14)
    p3.font.color.rgb = ACCENT_BLUE
    p3.font.name = 'Microsoft YaHei'
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(8)

add_textbox(slide, 1.5, 6.5, 10, 0.5, '覆盖从远视储备到早老花的全生命周期视觉管理', 16, MID_GRAY, True, PP_ALIGN.CENTER)

# ==================== P7: CDSA ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, 1.5, 1.0, 10, 0.6, 'P7', 14, MID_GRAY, False, PP_ALIGN.RIGHT)
add_textbox(slide, 1.5, 1.5, 10, 0.8, 'CDSA  ·  临床数据战略联盟', 32, DARK_BLUE, True)
add_textbox(slide, 1.5, 2.5, 10, 0.5, 'Clinical Data Strategic Alliance', 18, MID_GRAY, False)

# Flow boxes
flow = ['建档', '随访', 'AI辅助决策', '数据沉淀']
for i, step in enumerate(flow):
    left = 1.5 + i * 2.8
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(3.5), Inches(2.2), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE if i < 3 else ACCENT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    # Arrow
    if i < 3:
        add_textbox(slide, left + 2.2, 3.8, 0.6, 0.5, '→', 24, MID_GRAY, False, PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 5.2, 10, 0.5, '"效果到底怎么样？" ——数据沉淀下来，大家一起用', 16, DARK_GRAY, False, PP_ALIGN.CENTER)

# ==================== P8: 2026目标 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_textbox(slide, 1.5, 2.0, 10, 1.5, '5,000', 120, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 3.8, 10, 0.8, '例 标准化临床数据（2026）', 28, RGBColor(0xA0,0xC4,0xE8), False, PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 5.0, 10, 0.5, '病程管理软件  ·  免费部署', 20, RGBColor(0xCC,0xCC,0xCC), False, PP_ALIGN.CENTER)

# ==================== P9: 抗疲劳OK镜 - 故事 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, 1.5, 1.0, 10, 0.6, 'P9', 14, MID_GRAY, False, PP_ALIGN.RIGHT)
add_textbox(slide, 1.5, 1.8, 10, 0.8, '抗疲劳OK镜', 36, DARK_BLUE, True, PP_ALIGN.LEFT)
add_textbox(slide, 1.5, 2.6, 10, 0.5, '回到那个看不清的字', 20, MID_GRAY, False)

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10), Inches(3.0))
tf = txBox.text_frame
tf.word_wrap = True
lines = [
    '从业20年角塑医生 + 患者双重视角',
    '看手机、看处方、看小字 —— 近视力疲劳',
    '',
    '现有OK镜：核心参数都在解决近视控制和远视力',
    '看近的调节需求 —— 没有人专门做过',
]
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(22)
    p.font.color.rgb = DARK_GRAY
    p.font.name = 'Microsoft YaHei'
    p.space_after = Pt(6)

# ==================== P10: 抗疲劳OK镜 - 方案 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, 1.5, 1.0, 10, 0.6, 'P10', 14, MID_GRAY, False, PP_ALIGN.RIGHT)
add_textbox(slide, 1.5, 1.5, 10, 0.8, '方案与初步观察', 32, DARK_BLUE, True)

add_bullet_list(slide, 1.5, 2.8, 10, 3.5, [
    '在OK镜光学设计中加入近用调节参数调整',
    '本人亲身验证 + ~10位从业者/长期用户',
    '反馈一致：看近舒适度有改善',
    '',
    '⚠ 初步临床观察，非正式临床试验',
    '据我们了解，国际上尚属首次报告',
], 20, DARK_GRAY)

# ==================== P11: AI ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, 1.5, 1.0, 10, 0.6, 'P11', 14, MID_GRAY, False, PP_ALIGN.RIGHT)
add_textbox(slide, 1.5, 1.5, 10, 0.8, 'AI 与临床决策支持', 32, DARK_BLUE, True)

add_bullet_list(slide, 1.5, 2.8, 10, 3.0, [
    '辅助验配师做决策，不是替代验配师',
    '角膜地形图 + 屈光数据 → 自动推荐最优镜片方案',
    '历史随访数据 → 预测近视进展趋势',
    '异常指标自动预警',
    '',
    '联邦学习架构：数据不出本地，只交换模型参数',
], 20, DARK_GRAY)

# Demo placeholder
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(5.5), Inches(4), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_GRAY
shape.line.fill.background()
tf = shape.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
p = tf.paragraphs[0]
p.text = '[ Demo 现场演示 ]'
p.font.size = Pt(20)
p.font.color.rgb = MID_GRAY
p.font.name = 'Microsoft YaHei'

# ==================== P12: 总结 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_textbox(slide, 1.5, 1.5, 10, 0.8, '总结', 36, WHITE, True, PP_ALIGN.CENTER)

items = [
    '1.  已上市产品临床数据扎实，方向对',
    '2.  产品线扩展 —— 覆盖全生命周期',
    '3.  抗疲劳OK镜 —— 新尝试，期待与同行一起验证',
]
for i, item in enumerate(items):
    add_textbox(slide, 2.0, 2.8 + i * 0.9, 9, 0.7, item, 24, RGBColor(0xDD,0xDD,0xDD), False, PP_ALIGN.LEFT)

add_textbox(slide, 1.5, 6.0, 10, 0.5, '欢迎会后深入交流  ·  展台见', 18, RGBColor(0x99,0x99,0x99), False, PP_ALIGN.CENTER)

# Save
output_path = r'C:\Users\wangc\Downloads\ClarkWangVision\02-医学市场\4月9日产品发布会.pptx'
prs.save(output_path)
print(f'PPT saved to {output_path}')
