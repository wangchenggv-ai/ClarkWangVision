# -*- coding: utf-8 -*-
"""
研究背景与创新设计 - 3页PPT生成
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

PPT_FILE = r'C:\Users\wangc\Downloads\ClarkWangVision\离焦RGP_6个月临床数据分析报告.pptx'
OUTPUT_FILE = r'C:\Users\wangc\Downloads\ClarkWangVision\离焦RGP_6个月临床数据分析报告.pptx'

# 配色
COLORS = {
    'primary': RGBColor(0, 119, 182),
    'secondary': RGBColor(0, 180, 216),
    'accent': RGBColor(144, 224, 239),
    'dark': RGBColor(27, 42, 74),
    'white': RGBColor(255, 255, 255),
    'light_bg': RGBColor(248, 250, 252),
    'success': RGBColor(45, 155, 90),
    'warning': RGBColor(244, 162, 97),
    'danger': RGBColor(230, 57, 70),
    'gray': RGBColor(108, 117, 125),
    'light_blue': RGBColor(232, 245, 255),
}

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_bg_shape(slide):
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = COLORS['primary']; top_bar.line.fill.background()
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15))
    bottom_bar.fill.solid(); bottom_bar.fill.fore_color.rgb = COLORS['primary']; bottom_bar.line.fill.background()


def add_side_bar(slide):
    side_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT)
    side_bar.fill.solid(); side_bar.fill.fore_color.rgb = COLORS['primary']; side_bar.line.fill.background()


def add_page_number(slide, number, total):
    txBox = slide.shapes.add_textbox(Inches(12), Inches(7.1), Inches(1), Inches(0.3))
    p = txBox.text_frame.paragraphs[0]; p.text = f"{number}/{total}"; p.font.size = Pt(10); p.font.color.rgb = COLORS['gray']; p.alignment = PP_ALIGN.RIGHT


def add_footer_text(slide):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(4), Inches(0.3))
    p = txBox.text_frame.paragraphs[0]; p.text = '离焦RGP 临床数据分析'; p.font.size = Pt(9); p.font.color.rgb = COLORS['gray']


def add_title_box(slide, title, subtitle=None, y=Inches(0.3)):
    txBox = slide.shapes.add_textbox(Inches(0.5), y, Inches(12), Inches(0.7))
    p = txBox.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = COLORS['primary']
    if subtitle:
        p2 = txBox.text_frame.add_paragraph(); p2.text = subtitle; p2.font.size = Pt(14); p2.font.color.rgb = COLORS['gray']; p2.space_before = Pt(4)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y + Inches(0.75), Inches(2), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = COLORS['primary']; line.line.fill.background()


def add_rounded_card(slide, left, top, width, height, fill_color, border_color, text_lines):
    """添加圆角卡片并写入多行文字"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid(); card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color; card.line.width = Pt(2)
    tf = card.text_frame; tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
        p.space_before = Pt(4)
    return card


# ==================== 幻灯片 ====================

def create_slide_research_background(prs):
    """第1页：研究背景 — 既往RGP近视防控的困境"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide); add_side_bar(slide)
    add_title_box(slide, '研究背景', 'Research Background — 既往RGP近视防控的困境与挑战')

    # ---- 左栏：标准RGP的误区 ----
    add_rounded_card(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.8),
                     COLORS['light_bg'], COLORS['warning'],
                     [
                         ('标准RGP的"近视控制"神话已被打破', 16, True, COLORS['warning']),
                         ('', 6, False, COLORS['dark']),
                         ('长期以来，标准RGP被认为可通过物理压平角膜来控制近视', 12, False, COLORS['dark']),
                         ('', 4, False, COLORS['dark']),
                         ('CLAMP研究（Walline et al.）揭示真相：', 12, True, COLORS['dark']),
                         ('  - SER差异：RGP组 -1.56D vs 软镜组 -2.19D（p<0.001）', 11, False, COLORS['dark']),
                         ('  - 眼轴增长：RGP组 0.84mm vs 软镜组 0.79mm（p=0.57）', 11, False, COLORS['danger']),
                         ('  - 结论：标准RGP仅压平角膜掩盖度数，对眼轴增长无效', 12, True, COLORS['danger']),
                     ])

    # ---- 右栏：温州研究的失败 ----
    add_rounded_card(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.8),
                     COLORS['light_bg'], COLORS['danger'],
                     [
                         ('温州医科大学：离焦RGP对高度近视无效', 16, True, COLORS['danger']),
                         ('', 6, False, COLORS['dark']),
                         ('Yu et al. (2023) 纳入77名高度近视儿童（平均 -7.87D）', 12, False, COLORS['dark']),
                         ('', 4, False, COLORS['dark']),
                         ('2年随访结果：', 12, True, COLORS['dark']),
                         ('  - 1年眼轴增长：mRGP 0.20mm vs 对照 0.21mm（p=0.835）', 11, False, COLORS['dark']),
                         ('  - 2年眼轴增长：mRGP 0.37mm vs 对照 0.43mm（p=0.224）', 11, False, COLORS['dark']),
                         ('  - 结论：现有离焦设计对高度近视儿童无效', 12, True, COLORS['danger']),
                     ])

    # ---- 底部：问题核心 ----
    add_rounded_card(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.7),
                     COLORS['primary'], COLORS['primary'],
                     [
                         ('核心问题：我们缺少什么？', 18, True, COLORS['white']),
                         ('', 8, False, COLORS['white']),
                         ('1. 标准RGP无近视控制效果 — 物理压平角膜 ≠ 抑制眼轴增长', 14, False, COLORS['white']),
                         ('2. 现有离焦RGP对高度近视无效 — 光学离焦剂量无法跨越极端眼轴长度', 14, False, COLORS['white']),
                         ('3. 普通近视儿童（< -6D）的防控缺口 — OK镜有度数/散光限制，软镜矫正散光能力差', 14, False, COLORS['white']),
                         ('', 6, False, COLORS['white']),
                         ('我们需要一种全新的离焦RGP设计，针对普通近视儿童，提供有效且稳定的近视控制', 15, True, COLORS['accent']),
                     ])

    add_page_number(slide, 3, 13)
    add_footer_text(slide)


def create_slide_innovation(prs):
    """第2页：创新设计 — 本研究离焦RGP的独特之处"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide); add_side_bar(slide)
    add_title_box(slide, '创新设计', 'Innovation Design — 全新离焦RGP的光学突破')

    # ---- 左栏：设计原理 ----
    add_rounded_card(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(5.8),
                     COLORS['light_bg'], COLORS['primary'],
                     [
                         ('不是模拟OK镜 — 全新的光学架构', 16, True, COLORS['primary']),
                         ('', 6, False, COLORS['dark']),
                         ('传统离焦RGP：模拟OK镜的逆几何设计，离焦区在中周部', 12, False, COLORS['gray']),
                         ('本研究设计：光学作用聚焦于瞳孔区，全新思路', 12, True, COLORS['dark']),
                         ('', 8, False, COLORS['dark']),
                         ('三大核心创新：', 14, True, COLORS['primary']),
                         ('', 4, False, COLORS['dark']),
                         ('① 2mm瞳孔区精准作用', 13, True, COLORS['dark']),
                         ('   光学离焦信号直接作用于瞳孔核心区域', 11, False, COLORS['gray']),
                         ('   日间瞳孔收缩时仍能有效激活，避免"光学休眠"', 11, False, COLORS['gray']),
                         ('', 4, False, COLORS['dark']),
                         ('② 离焦环稳定2-3D', 13, True, COLORS['dark']),
                         ('   不随瞬目（眨眼）滑动而衰减', 11, False, COLORS['gray']),
                         ('   全天候持续输出稳定的近视离焦信号', 11, False, COLORS['gray']),
                         ('', 4, False, COLORS['dark']),
                         ('③ 不影响中心视力', 13, True, COLORS['dark']),
                         ('   精确的光学区设计，保证黄斑中心凹成像质量', 11, False, COLORS['gray']),
                         ('   兼顾近视控制与清晰视力，学习生活不受影响', 11, False, COLORS['gray']),
                     ])

    # ---- 右栏：入组设计与研究目标 ----
    add_rounded_card(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.6),
                     COLORS['light_bg'], COLORS['success'],
                     [
                         ('入组设计：排除高度近视干扰', 16, True, COLORS['success']),
                         ('', 6, False, COLORS['dark']),
                         ('入组标准：SER < -6.00D 的儿童', 13, True, COLORS['dark']),
                         ('', 4, False, COLORS['dark']),
                         ('设计考量：', 12, True, COLORS['dark']),
                         ('  - 高度近视（≥ -6D）眼球后极部严重变形', 11, False, COLORS['gray']),
                         ('  - 周边视网膜极度后退，离焦信号难以到达', 11, False, COLORS['gray']),
                         ('  - 排除这一混杂因素，聚焦普通近视儿童', 11, False, COLORS['gray']),
                     ])

    add_rounded_card(slide, Inches(6.8), Inches(4.2), Inches(6), Inches(2.9),
                     COLORS['primary'], COLORS['primary'],
                     [
                         ('研究目标', 16, True, COLORS['white']),
                         ('', 6, False, COLORS['white']),
                         ('首个针对普通近视儿童的创新离焦RGP研究', 14, True, COLORS['accent']),
                         ('', 4, False, COLORS['white']),
                         ('核心科学问题：', 13, True, COLORS['white']),
                         ('这种作用于瞳孔区、离焦环稳定、不影响视力的', 12, False, COLORS['white']),
                         ('全新离焦RGP设计，能否有效控制普通近视儿童的眼轴增长？', 12, False, COLORS['accent']),
                     ])

    add_page_number(slide, 4, 13)
    add_footer_text(slide)


def create_slide_significance(prs):
    """第3页：研究意义与临床价值"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide); add_side_bar(slide)
    add_title_box(slide, '研究意义', 'Clinical Significance — 填补近视防控领域的关键空白')

    # ---- 顶部：现有防控手段对比 ----
    add_rounded_card(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.8),
                     COLORS['light_bg'], COLORS['secondary'],
                     [
                         ('现有近视防控光学手段的局限性', 16, True, COLORS['secondary']),
                         ('', 4, False, COLORS['dark']),
                         ('OK镜：度数 ≤ -6D、散光 ≤ -1.50D 限制    |    多焦点软镜：散光 > 1.50D 时视力严重下降    |    离焦框架镜：眼球转动时离焦信号被破坏', 12, False, COLORS['dark']),
                         ('标准RGP：无近视控制效果    |    既往离焦RGP：高度近视无效    |    缺口：普通近视儿童缺乏有效的RGP类防控方案', 12, True, COLORS['danger']),
                     ])

    # ---- 中间三栏 ----
    # 栏1
    add_rounded_card(slide, Inches(0.5), Inches(3.4), Inches(3.9), Inches(3.6),
                     COLORS['success'], COLORS['success'],
                     [
                         ('学术价值', 16, True, COLORS['white']),
                         ('', 8, False, COLORS['white']),
                         ('首个创新设计的离焦RGP', 13, True, COLORS['accent']),
                         ('临床研究', 13, True, COLORS['accent']),
                         ('', 6, False, COLORS['white']),
                         ('区别于模拟OK镜的传统设计', 11, False, COLORS['white']),
                         ('验证瞳孔区作用的全新', 11, False, COLORS['white']),
                         ('光学离焦理论', 11, False, COLORS['white']),
                         ('', 6, False, COLORS['white']),
                         ('填补普通近视儿童', 11, False, COLORS['white']),
                         ('RGP防控方案的空白', 11, False, COLORS['white']),
                     ])

    # 栏2
    add_rounded_card(slide, Inches(4.7), Inches(3.4), Inches(3.9), Inches(3.6),
                     COLORS['primary'], COLORS['primary'],
                     [
                         ('临床价值', 16, True, COLORS['white']),
                         ('', 8, False, COLORS['white']),
                         ('为 < -6D 儿童提供', 13, True, COLORS['accent']),
                         ('新的防控选择', 13, True, COLORS['accent']),
                         ('', 6, False, COLORS['white']),
                         ('保留RGP矫正散光的优势', 11, False, COLORS['white']),
                         ('同时实现近视离焦控制', 11, False, COLORS['white']),
                         ('', 6, False, COLORS['white']),
                         ('日间佩戴，感染风险', 11, False, COLORS['white']),
                         ('低于夜戴型OK镜', 11, False, COLORS['white']),
                     ])

    # 栏3
    add_rounded_card(slide, Inches(8.9), Inches(3.4), Inches(3.9), Inches(3.6),
                     COLORS['warning'], COLORS['warning'],
                     [
                         ('未来方向', 16, True, COLORS['white']),
                         ('', 8, False, COLORS['white']),
                         ('如6个月数据提示有效', 13, True, COLORS['accent']),
                         ('将开展：', 13, True, COLORS['accent']),
                         ('', 6, False, COLORS['white']),
                         ('- 延长随访至12-24个月', 11, False, COLORS['white']),
                         ('- 增设对照组验证', 11, False, COLORS['white']),
                         ('- 扩大样本量', 11, False, COLORS['white']),
                         ('- 探索联合低浓度阿托品', 11, False, COLORS['white']),
                         ('  对控制欠佳者的增效方案', 11, False, COLORS['white']),
                     ])

    add_page_number(slide, 5, 13)
    add_footer_text(slide)


# ==================== 主程序 ====================
def main():
    # 加载现有PPT
    prs = Presentation(PPT_FILE)

    # 获取当前页数
    current_slides = len(prs.slides)
    total_slides = current_slides + 3

    print(f'当前PPT页数: {current_slides}')
    print('开始添加3页背景与创新设计幻灯片...')

    # 在目录页后、研究概况页前插入（索引2的位置）
    # 但由于python-pptx不支持直接在中间插入，我们在末尾添加，用户可手动调整顺序
    create_slide_research_background(prs)
    print(f'  [{current_slides+1}/{total_slides}] 研究背景')

    create_slide_innovation(prs)
    print(f'  [{current_slides+2}/{total_slides}] 创新设计')

    create_slide_significance(prs)
    print(f'  [{current_slides+3}/{total_slides}] 研究意义')

    prs.save(OUTPUT_FILE)
    print(f'\nPPT已更新保存: {OUTPUT_FILE}')
    print('新增3页位于PPT末尾，请在PowerPoint中移动到目录页之后')


if __name__ == '__main__':
    main()
