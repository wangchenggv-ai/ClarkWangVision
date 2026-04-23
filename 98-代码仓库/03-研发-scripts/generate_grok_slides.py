# -*- coding: utf-8 -*-
"""
基于Grok建议的3页背景PPT生成
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_FILE = r'C:\Users\wangc\Downloads\ClarkWangVision\离焦RGP_研究背景与创新设计.pptx'

# 配色
C = {
    'pri': RGBColor(0, 119, 182),
    'sec': RGBColor(0, 180, 216),
    'acc': RGBColor(144, 224, 239),
    'dark': RGBColor(27, 42, 74),
    'white': RGBColor(255, 255, 255),
    'light': RGBColor(248, 250, 252),
    'green': RGBColor(45, 155, 90),
    'orange': RGBColor(244, 162, 97),
    'red': RGBColor(230, 57, 70),
    'gray': RGBColor(108, 117, 125),
}

W = Inches(13.333)
H = Inches(7.5)


def add_bg(slide):
    for y, h in [(0, 0.08), (7.35, 0.15)]:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(y), W, Inches(h))
        bar.fill.solid(); bar.fill.fore_color.rgb = C['pri']; bar.line.fill.background()
    side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), H)
    side.fill.solid(); side.fill.fore_color.rgb = C['pri']; side.line.fill.background()


def add_title(slide, title, sub, y=0.3):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12), Inches(0.7))
    p = txBox.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = C['pri']
    if sub:
        p2 = txBox.text_frame.add_paragraph()
        p2.text = sub; p2.font.size = Pt(14); p2.font.color.rgb = C['gray']; p2.space_before = Pt(4)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y) + Inches(0.75), Inches(2), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = C['pri']; line.line.fill.background()


def card(slide, l, t, w, h, fill, border, lines):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = Pt(2)
    tf = shp.text_frame; tf.word_wrap = True
    for i, (txt, sz, bold, clr) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = clr
        p.space_before = Pt(3)
    return shp


def footer(slide, n, total):
    txBox = slide.shapes.add_textbox(Inches(12), Inches(7.1), Inches(1), Inches(0.3))
    p = txBox.text_frame.paragraphs[0]; p.text = f"{n}/{total}"; p.font.size = Pt(10); p.font.color.rgb = C['gray']; p.alignment = PP_ALIGN.RIGHT
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(4), Inches(0.3))
    p2 = txBox2.text_frame.paragraphs[0]; p2.text = '离焦RGP 临床数据分析'; p2.font.size = Pt(9); p2.font.color.rgb = C['gray']


# ==================== 幻灯片 ====================

def slide1(prs):
    """幻灯片1：背景与既往研究的局限性"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, 'RGP近视防控研究的背景', '标准RGP神话破灭与高度近视离焦RGP的局限')

    # 左栏：标准RGP
    card(slide, 0.5, 1.3, 6, 2.5, C['light'], C['orange'], [
        ('标准RGP的"近视控制"神话已被打破', 15, True, C['orange']),
        ('', 4, False, C['dark']),
        ('CLAMP研究（Walline等，3年RCT）：', 12, True, C['dark']),
        ('  SER差异：-1.56D vs -2.19D（p<0.001）', 11, False, C['dark']),
        ('  眼轴延长：0.84mm vs 0.79mm（p=0.57）', 11, False, C['red']),
        ('  结论：仅机械压平角膜产生SER假象', 11, False, C['dark']),
        ('  无法调控视网膜-巩膜生长信号', 11, True, C['red']),
    ])

    # 右栏：温州研究
    card(slide, 6.8, 1.3, 6, 2.5, C['light'], C['red'], [
        ('温州医科大学：离焦RGP对高度近视无效', 15, True, C['red']),
        ('', 4, False, C['dark']),
        ('Yu et al. (2023)，高度近视儿童（平均-7.87D）：', 12, True, C['dark']),
        ('  1年：0.20±0.17mm vs 0.21±0.14mm（p=0.835）', 11, False, C['dark']),
        ('  2年：0.37±0.27mm vs 0.43±0.23mm（p=0.224）', 11, False, C['dark']),
        ('  原因：极度长椭圆形眼球→离焦"剂量"不足', 11, True, C['red']),
        ('  中央区设计过大→日间瞳孔收缩时光学休眠', 11, False, C['gray']),
    ])

    # 底部：关键数据对比表
    card(slide, 0.5, 4.1, 12.3, 1.3, C['pri'], C['pri'], [
        ('关键证据对比', 14, True, C['white']),
        ('CLAMP研究：标准RGP → 眼轴无差异（p=0.57）    |    温州研究：离焦RGP对高度近视 → 眼轴无差异（p=0.835/0.224）', 11, False, C['white']),
        ('结论：既往设计均无法有效抑制眼轴增长，需开发全新光学方案', 12, True, C['acc']),
    ])

    # 底部结论
    card(slide, 0.5, 5.6, 12.3, 1.3, C['light'], C['sec'], [
        ('核心洞察：既往设计局限凸显需求', 15, True, C['sec']),
        ('', 4, False, C['dark']),
        ('近视防控核心目标：将眼轴年增长量压制至生理基线水平（≤0.10–0.15 mm/年）', 12, False, C['dark']),
        ('既往RGP方案均未达成此目标 → 需针对普通低中度近视儿童开发全新光学方案', 12, True, C['pri']),
    ])

    footer(slide, 1, 3)


def slide2(prs):
    """幻灯片2：本研究离焦RGP的创新设计"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, '本研究创新设计', '首款针对2mm瞳孔区的稳定离焦RGP设计')

    # 左栏：核心创新
    card(slide, 0.5, 1.3, 6, 5.7, C['light'], C['pri'], [
        ('核心光学创新：三大突破', 16, True, C['pri']),
        ('', 6, False, C['dark']),
        ('突破一：2mm瞳孔区精准作用', 13, True, C['dark']),
        ('  光学离焦信号直接作用于瞳孔核心区域', 11, False, C['gray']),
        ('  日间瞳孔收缩时仍能有效激活，避免"光学休眠"', 11, False, C['gray']),
        ('  区别于传统mRGP中周部离焦设计', 11, False, C['gray']),
        ('', 4, False, C['dark']),
        ('突破二：离焦环稳定2–3D', 13, True, C['dark']),
        ('  不随瞬目（眨眼）滑动而衰减', 11, False, C['gray']),
        ('  全天候持续输出稳定的近视离焦信号', 11, False, C['gray']),
        ('  最大化激活周边视网膜"停止生长"信号', 11, False, C['gray']),
        ('', 4, False, C['dark']),
        ('突破三：不影响中心视力', 13, True, C['dark']),
        ('  精确光学区设计，保证黄斑中心凹成像质量', 11, False, C['gray']),
        ('  泪液透镜确保卓越中心成像，消除模糊生长信号', 11, False, C['gray']),
        ('  兼顾近视控制与清晰视力', 11, False, C['gray']),
    ])

    # 右栏：设计优势对比
    card(slide, 6.8, 1.3, 6, 2.5, C['light'], C['green'], [
        ('设计优势：突破传统模式', 15, True, C['green']),
        ('', 4, False, C['dark']),
        ('非Ortho-K逆几何压痕（夜间佩戴）', 12, False, C['dark']),
        ('非软镜包裹复制角膜地形（散光矫正差）', 12, False, C['dark']),
        ('独立优化日间稳定离焦轮廓', 12, True, C['dark']),
        ('RGP材质高透氧性（Dk>80）+ 安全日间佩戴', 12, False, C['green']),
    ])

    card(slide, 6.8, 4.1, 6, 2.9, C['pri'], C['pri'], [
        ('与现有手段的本质差异', 15, True, C['white']),
        ('', 4, False, C['white']),
        ('传统离焦RGP：模拟OK镜，离焦区在中周部', 11, False, C['acc']),
        ('本研究设计：光学作用聚焦于瞳孔区，全新架构', 11, True, C['white']),
        ('', 4, False, C['white']),
        ('避免既往mRGP在高度近视中的光学衰减：', 11, False, C['white']),
        ('- 视网膜非对称性问题', 10, False, C['acc']),
        ('- 镜片偏心导致的离焦信号不均匀', 10, False, C['acc']),
        ('- 中央区过大导致的光学休眠', 10, False, C['acc']),
    ])

    footer(slide, 2, 3)


def slide3(prs):
    """幻灯片3：临床试验设计初衷与目标人群"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, '临床研究设计初衷', '验证新型离焦RGP在普通儿童中的近视防控效果')

    # 左栏：研究逻辑5点
    card(slide, 0.5, 1.3, 6, 5.7, C['light'], C['pri'], [
        ('研究设计逻辑（5点对应）', 16, True, C['pri']),
        ('', 6, False, C['dark']),
        ('① 既往认知：标准RGP有防控效果 → 实际无效果', 12, True, C['orange']),
        ('   CLAMP研究证实仅角膜压平假象，眼轴无效', 11, False, C['gray']),
        ('', 3, False, C['dark']),
        ('② 既往尝试：温州离焦RGP对高度近视 → 无效', 12, True, C['red']),
        ('   极端眼轴导致离焦剂量不足，光学休眠', 11, False, C['gray']),
        ('', 3, False, C['dark']),
        ('③ 本研究创新：全新设计，非模拟OK镜', 12, True, C['green']),
        ('   作用于2mm瞳孔区，离焦环稳定2-3D，不影响视力', 11, False, C['gray']),
        ('', 3, False, C['dark']),
        ('④ 入组设计：< -6D，排除高度近视干扰', 12, True, C['pri']),
        ('   聚焦普通近视儿童，避免极端眼轴的混杂因素', 11, False, C['gray']),
        ('', 3, False, C['dark']),
        ('⑤ 研究目标：证明普通儿童中此设计的防控效果', 12, True, C['sec']),
        ('   填补现有证据空白', 11, False, C['gray']),
    ])

    # 右栏：目标与意义
    card(slide, 6.8, 1.3, 6, 2.5, C['light'], C['green'], [
        ('研究初衷与目标', 15, True, C['green']),
        ('', 4, False, C['dark']),
        ('目标人群：普通低中度近视儿童（< -6D）', 13, True, C['dark']),
        ('', 4, False, C['dark']),
        ('核心科学问题：', 12, True, C['dark']),
        ('优化后的离焦RGP能否将眼轴年增长量', 11, False, C['dark']),
        ('有效压制至生理基线（≤0.10–0.15 mm/年）', 11, True, C['green']),
    ])

    card(slide, 6.8, 4.1, 6, 1.5, C['light'], C['sec'], [
        ('主要终点指标', 15, True, C['sec']),
        ('', 4, False, C['dark']),
        ('眼轴长度（AL）| 等效球镜（SER）', 12, True, C['dark']),
        ('安全性（角膜染色、感染率）| 依从性 | 视觉质量', 11, False, C['gray']),
    ])

    # 底部结论框
    card(slide, 0.5, 5.6, 12.3, 1.3, C['pri'], C['pri'], [
        ('预期临床意义', 14, True, C['white']),
        ('', 4, False, C['white']),
        ('为低中度近视儿童提供安全、有效、日间佩戴的创新选择', 12, False, C['white']),
        ('确立新型离焦RGP在近视防控武器库中的定位 — 本设计直接回应既往局限，开启普通儿童精准防控新路径', 12, True, C['acc']),
    ])

    footer(slide, 3, 3)


# ==================== 主程序 ====================
def main():
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H

    print('开始生成3页背景PPT...')
    slide1(prs); print('  [1/3] 背景与既往研究局限性')
    slide2(prs); print('  [2/3] 创新设计')
    slide3(prs); print('  [3/3] 研究初衷与目标人群')

    prs.save(OUTPUT_FILE)
    print(f'\nPPT已保存: {OUTPUT_FILE}')


if __name__ == '__main__':
    main()
