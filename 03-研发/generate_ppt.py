# -*- coding: utf-8 -*-
"""
离焦RGP 6个月临床数据分析 - PPT自动生成器
深蓝医疗风格
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ==================== 配置 ====================
CHARTS_DIR = r'C:\Users\wangc\Downloads\ClarkWangVision\charts'
OUTPUT_FILE = r'C:\Users\wangc\Downloads\ClarkWangVision\离焦RGP_6个月临床数据分析报告.pptx'

# 深蓝医疗配色
COLORS = {
    'primary': RGBColor(0, 119, 182),      # #0077B6 深蓝
    'secondary': RGBColor(0, 180, 216),     # #00B4D8 浅蓝
    'accent': RGBColor(144, 224, 239),      # #90E0EF 天蓝
    'dark': RGBColor(27, 42, 74),           # #1B2A4A 深色文字
    'white': RGBColor(255, 255, 255),
    'light_bg': RGBColor(248, 250, 252),    # #F8FAFC 浅背景
    'success': RGBColor(45, 155, 90),       # #2D9B5A 绿色
    'warning': RGBColor(244, 162, 97),      # #F4A261 橙色
    'danger': RGBColor(230, 57, 70),        # #E63946 红色
    'gray': RGBColor(108, 117, 125),        # #6C757D 灰色
}

# 页面尺寸（16:9）
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_bg_shape(slide, color=COLORS['primary']):
    """添加背景色块"""
    # 顶部装饰条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = color
    top_bar.line.fill.background()
    
    # 底部装饰条
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = color
    bottom_bar.line.fill.background()


def add_side_bar(slide, color=COLORS['primary']):
    """添加左侧装饰条"""
    side_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT)
    side_bar.fill.solid()
    side_bar.fill.fore_color.rgb = color
    side_bar.line.fill.background()


def add_page_number(slide, number, total):
    """添加页码"""
    txBox = slide.shapes.add_textbox(Inches(12), Inches(7.1), Inches(1), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{number}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS['gray']
    p.alignment = PP_ALIGN.RIGHT


def add_footer_text(slide, text='离焦RGP 临床数据分析'):
    """添加页脚文字"""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(4), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS['gray']


def add_title_box(slide, title, subtitle=None, y=Inches(0.3)):
    """添加页面标题"""
    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), y, Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 副标题
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLORS['gray']
        p2.space_before = Pt(4)
    
    # 标题下划线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), y + Inches(0.75), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['primary']
    line.line.fill.background()


# ==================== 幻灯片创建 ====================

def create_cover_slide(prs):
    """封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 全屏深蓝背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()
    
    # 装饰圆
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9.5), Inches(-1), Inches(5), Inches(5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS['secondary']
    circle.line.fill.background()
    
    # 装饰圆2
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(-2), Inches(4), Inches(4), Inches(4))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = COLORS['accent']
    circle2.line.fill.background()
    
    # 主标题
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '离焦RGP 6个月临床数据分析报告'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.LEFT
    
    # 副标题
    p2 = tf.add_paragraph()
    p2.text = 'Defocus RGP Clinical Data Analysis Report'
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLORS['accent']
    p2.space_before = Pt(12)
    
    # 分隔线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.8), Inches(3), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent']
    line.line.fill.background()
    
    # 研究信息
    info_items = [
        '样本量：11人 / 22眼',
        '随访周期：6个月',
        '平均年龄：10.5 ± 1.7 岁',
        '年化眼轴增长：+0.181mm（控制效力约50~55%）'
    ]
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(2.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(info_items):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f'▸  {item}'
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['white']
        p.space_before = Pt(8)
    
    # 页脚
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = '内部讨论使用  |  未经同行评审'
    p3.font.size = Pt(12)
    p3.font.color.rgb = COLORS['accent']
    p3.font.italic = True


def create_toc_slide(prs):
    """目录页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide)
    
    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = '目录'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 英文副标题
    p2 = tf.add_paragraph()
    p2.text = 'CONTENTS'
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLORS['gray']
    p2.font.bold = True
    
    # 目录项
    toc_items = [
        ('01', '研究概况', '样本特征与基线数据'),
        ('02', '核心疗效指标', '眼轴增长、屈光度变化、矫正视力'),
        ('03', '亚组分析', '年龄分层、屈光度分层'),
        ('04', '个体差异与关注对象', '良好控制与控制欠佳案例'),
        ('05', '与文献基准对比', '控制效力评估'),
        ('06', '满意度评价', '戴镜舒适度与视力质量'),
        ('07', '局限性与建议', '研究不足与后续计划'),
    ]
    
    for i, (num, title, desc) in enumerate(toc_items):
        y = Inches(1.8) + Inches(i * 0.7)
        
        # 编号
        num_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), y, Inches(0.6), Inches(0.5))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = COLORS['primary']
        num_box.line.fill.background()
        ntf = num_box.text_frame
        ntf.paragraphs[0].text = num
        ntf.paragraphs[0].font.size = Pt(16)
        ntf.paragraphs[0].font.bold = True
        ntf.paragraphs[0].font.color.rgb = COLORS['white']
        ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(1.8), y, Inches(4), Inches(0.3))
        ttf = title_box.text_frame
        ttf.paragraphs[0].text = title
        ttf.paragraphs[0].font.size = Pt(18)
        ttf.paragraphs[0].font.bold = True
        ttf.paragraphs[0].font.color.rgb = COLORS['dark']
        
        # 描述
        desc_box = slide.shapes.add_textbox(Inches(1.8), y + Inches(0.28), Inches(5), Inches(0.25))
        dtf = desc_box.text_frame
        dtf.paragraphs[0].text = desc
        dtf.paragraphs[0].font.size = Pt(11)
        dtf.paragraphs[0].font.color.rgb = COLORS['gray']
    
    # 右侧装饰图
    deco = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9), Inches(2), Inches(3.5), Inches(3.5))
    deco.fill.solid()
    deco.fill.fore_color.rgb = COLORS['accent']
    deco.line.fill.background()
    
    add_page_number(slide, 2, 10)
    add_footer_text(slide)


def create_overview_slide(prs):
    """研究概况页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide)
    add_side_bar(slide)
    add_title_box(slide, '研究概况', 'Study Overview')
    
    # 插入图表
    chart_path = os.path.join(CHARTS_DIR, '08_研究概况.png')
    if os.path.exists(chart_path):
        slide.shapes.add_picture(chart_path, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5))
    
    add_page_number(slide, 3, 10)
    add_footer_text(slide)


def create_al_trend_slide(prs):
    """眼轴增长趋势页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide)
    add_side_bar(slide)
    add_title_box(slide, '眼轴增长控制趋势', 'Axial Length Change Trend')
    
    chart_path = os.path.join(CHARTS_DIR, '01_眼轴增长趋势.png')
    if os.path.exists(chart_path):
        slide.shapes.add_picture(chart_path, Inches(1), Inches(1.3), Inches(10), Inches(5.8))
    
    # 关键结论
    key_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(5.5), Inches(4.3), Inches(1.5))
    key_box.fill.solid()
    key_box.fill.fore_color.rgb = RGBColor(232, 245, 233)
    key_box.line.color.rgb = COLORS['success']
    
    ktf = key_box.text_frame
    ktf.word_wrap = True
    ktf.paragraphs[0].text = '核心结论'
    ktf.paragraphs[0].font.size = Pt(14)
    ktf.paragraphs[0].font.bold = True
    ktf.paragraphs[0].font.color.rgb = COLORS['success']
    
    kp = ktf.add_paragraph()
    kp.text = '年化眼轴增长 +0.181mm'
    kp.font.size = Pt(13)
    kp.font.bold = True
    kp.font.color.rgb = COLORS['dark']
    
    kp2 = ktf.add_paragraph()
    kp2.text = '眼轴控制效力约50~55%'
    kp2.font.size = Pt(12)
    kp2.font.color.rgb = COLORS['dark']
    
    kp3 = ktf.add_paragraph()
    kp3.text = '与离焦软镜处于同一效力量级'
    kp3.font.size = Pt(12)
    kp3.font.color.rgb = COLORS['gray']
    
    add_page_number(slide, 4, 10)
    add_footer_text(slide)


def create_ser_slide(prs):
    """SER变化分类页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide)
    add_side_bar(slide)
    add_title_box(slide, '屈光度(SER)变化分析', 'Spherical Equivalent Refraction Change')
    
    chart_path = os.path.join(CHARTS_DIR, '02_SER变化分类.png')
    if os.path.exists(chart_path):
        slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.2), Inches(7), Inches(6))
    
    # 右侧数据卡片
    cards = [
        ('稳定', '10眼', '45%', COLORS['success']),
        ('进展', '7眼', '32%', COLORS['warning']),
        ('改善', '5眼', '23%', COLORS['primary']),
    ]
    
    for i, (label, count, pct, color) in enumerate(cards):
        y = Inches(1.5) + Inches(i * 1.8)
        
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), y, Inches(4), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['white']
        card.line.color.rgb = color
        card.line.width = Pt(2)
        
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.paragraphs[0].text = label
        ctf.paragraphs[0].font.size = Pt(16)
        ctf.paragraphs[0].font.bold = True
        ctf.paragraphs[0].font.color.rgb = color
        
        cp = ctf.add_paragraph()
        cp.text = f'{count}  ({pct})'
        cp.font.size = Pt(24)
        cp.font.bold = True
        cp.font.color.rgb = COLORS['dark']
    
    # 平均值卡片
    avg_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(6), Inches(4), Inches(1))
    avg_card.fill.solid()
    avg_card.fill.fore_color.rgb = COLORS['primary']
    avg_card.line.fill.background()
    
    atf = avg_card.text_frame
    atf.paragraphs[0].text = '平均变化：-0.10 ± 0.34 D'
    atf.paragraphs[0].font.size = Pt(16)
    atf.paragraphs[0].font.bold = True
    atf.paragraphs[0].font.color.rgb = COLORS['white']
    atf.paragraphs[0].alignment = PP_ALIGN.CENTER
    atf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    add_page_number(slide, 5, 10)
    add_footer_text(slide)


def create_subgroup_slide(prs):
    """亚组分析页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide)
    add_side_bar(slide)
    add_title_box(slide, '亚组分析', 'Subgroup Analysis')
    
    chart_age = os.path.join(CHARTS_DIR, '03_年龄分层分析.png')
    chart_ser = os.path.join(CHARTS_DIR, '04_屈光度分层分析.png')
    
    if os.path.exists(chart_age):
        slide.shapes.add_picture(chart_age, Inches(0.3), Inches(1.3), Inches(6.2), Inches(5.5))
    
    if os.path.exists(chart_ser):
        slide.shapes.add_picture(chart_ser, Inches(6.8), Inches(1.3), Inches(6.2), Inches(5.5))
    
    # 分析说明
    note_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.8))
    note_box.fill.solid()
    note_box.fill.fore_color.rgb = COLORS['light_bg']
    note_box.line.color.rgb = COLORS['accent']
    
    ntf = note_box.text_frame
    ntf.word_wrap = True
    ntf.paragraphs[0].text = '▸ 低龄组（≤8岁）控制效果更优，但样本量小（n=6），需谨慎解读    ▸ 中高度近视者眼轴增长更快，符合临床预期'
    ntf.paragraphs[0].font.size = Pt(12)
    ntf.paragraphs[0].font.color.rgb = COLORS['dark']
    ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    add_page_number(slide, 6, 10)
    add_footer_text(slide)


def create_control_slide(prs):
    """疗效控制率页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide)
    add_side_bar(slide)
    add_title_box(slide, '个体疗效分布', 'Individual Efficacy Distribution')
    
    chart_path = os.path.join(CHARTS_DIR, '05_疗效控制率.png')
    if os.path.exists(chart_path):
        slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.3), Inches(12), Inches(4))
    
    # 关注对象卡片
    warn_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.5), Inches(6), Inches(1.5))
    warn_card.fill.solid()
    warn_card.fill.fore_color.rgb = RGBColor(255, 243, 224)
    warn_card.line.color.rgb = COLORS['warning']
    
    wtf = warn_card.text_frame
    wtf.word_wrap = True
    wtf.paragraphs[0].text = '控制欠佳关注对象'
    wtf.paragraphs[0].font.size = Pt(14)
    wtf.paragraphs[0].font.bold = True
    wtf.paragraphs[0].font.color.rgb = COLORS['warning']
    
    wp = wtf.add_paragraph()
    wp.text = 'WMZH（OD +0.33mm）、ZJHA-0011（OS +0.21mm）'
    wp.font.size = Pt(12)
    wp.font.color.rgb = COLORS['dark']
    
    wp2 = wtf.add_paragraph()
    wp2.text = '均为11岁男性，需排查依从性、用眼习惯等因素'
    wp2.font.size = Pt(11)
    wp2.font.color.rgb = COLORS['gray']
    
    # 良好控制卡片
    good_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(5.5), Inches(6), Inches(1.5))
    good_card.fill.solid()
    good_card.fill.fore_color.rgb = RGBColor(232, 245, 233)
    good_card.line.color.rgb = COLORS['success']
    
    gtf = good_card.text_frame
    gtf.word_wrap = True
    gtf.paragraphs[0].text = '良好控制案例'
    gtf.paragraphs[0].font.size = Pt(14)
    gtf.paragraphs[0].font.bold = True
    gtf.paragraphs[0].font.color.rgb = COLORS['success']
    
    gp = gtf.add_paragraph()
    gp.text = 'XYRA 双眼眼轴缩短、YUGU 双眼眼轴零增长'
    gp.font.size = Pt(12)
    gp.font.color.rgb = COLORS['dark']
    
    gp2 = gtf.add_paragraph()
    gp2.text = '中位数 +0.07mm | Q1=+0.02mm | Q3=+0.18mm'
    gp2.font.size = Pt(11)
    gp2.font.color.rgb = COLORS['gray']
    
    add_page_number(slide, 7, 10)
    add_footer_text(slide)


def create_comparison_slide(prs):
    """文献基准对比页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide)
    add_side_bar(slide)
    add_title_box(slide, '与文献基准对比', 'Comparison with Literature Baseline')
    
    chart_path = os.path.join(CHARTS_DIR, '07_文献基准对比.png')
    if os.path.exists(chart_path):
        slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.3), Inches(10), Inches(5.8))
    
    # 结论卡片
    conclusion = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(3), Inches(4.3), Inches(3.5))
    conclusion.fill.solid()
    conclusion.fill.fore_color.rgb = COLORS['primary']
    conclusion.line.fill.background()
    
    ctf = conclusion.text_frame
    ctf.word_wrap = True
    ctf.paragraphs[0].text = '结论'
    ctf.paragraphs[0].font.size = Pt(20)
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].font.color.rgb = COLORS['white']
    ctf.paragraphs[0].space_after = Pt(12)
    
    conclusions = [
        '年化眼轴增长 0.181mm',
        ' vs 未干预基准 0.30~0.40mm',
        '',
        '眼轴控制效力',
        '约50~55%',
        '',
        '与离焦软镜（DIMS/DISC）',
        '处于同一效力量级'
    ]
    
    for text in conclusions:
        cp = ctf.add_paragraph()
        cp.text = text
        if text.startswith('年化') or text.startswith('约') or text.startswith('处于'):
            cp.font.size = Pt(14)
            cp.font.bold = True
            cp.font.color.rgb = COLORS['accent']
        elif text.startswith('vs') or text.startswith('眼轴'):
            cp.font.size = Pt(12)
            cp.font.color.rgb = COLORS['white']
        else:
            cp.font.size = Pt(11)
            cp.font.color.rgb = RGBColor(200, 220, 240)
        cp.space_before = Pt(2)
    
    add_page_number(slide, 8, 10)
    add_footer_text(slide)


def create_satisfaction_slide(prs):
    """满意度页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide)
    add_side_bar(slide)
    add_title_box(slide, '患者满意度评价', 'Patient Satisfaction Assessment')
    
    chart_path = os.path.join(CHARTS_DIR, '06_满意度趋势.png')
    if os.path.exists(chart_path):
        slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.3), Inches(8), Inches(5.5))
    
    # 维度卡片
    dimensions = ['远视力', '近视力', '行走舒适度', '夜间视力']
    for i, dim in enumerate(dimensions):
        y = Inches(1.5) + Inches(i * 1.3)
        
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9), y, Inches(3.8), Inches(1.1))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['white']
        card.line.color.rgb = COLORS['secondary']
        card.line.width = Pt(1.5)
        
        ctf = card.text_frame
        ctf.paragraphs[0].text = dim
        ctf.paragraphs[0].font.size = Pt(14)
        ctf.paragraphs[0].font.bold = True
        ctf.paragraphs[0].font.color.rgb = COLORS['primary']
        
        cp = ctf.add_paragraph()
        cp.text = '表现良好'
        cp.font.size = Pt(11)
        cp.font.color.rgb = COLORS['success']
        cp.font.bold = True
    
    add_page_number(slide, 9, 10)
    add_footer_text(slide)


def create_limitation_slide(prs):
    """局限性与建议页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide)
    add_side_bar(slide)
    add_title_box(slide, '局限性与后续建议', 'Limitations & Recommendations')
    
    # 局限性卡片
    limit_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(6), Inches(5.5))
    limit_card.fill.solid()
    limit_card.fill.fore_color.rgb = RGBColor(255, 243, 224)
    limit_card.line.color.rgb = COLORS['warning']
    
    ltf = limit_card.text_frame
    ltf.word_wrap = True
    ltf.paragraphs[0].text = '研究局限性'
    ltf.paragraphs[0].font.size = Pt(20)
    ltf.paragraphs[0].font.bold = True
    ltf.paragraphs[0].font.color.rgb = COLORS['warning']
    ltf.paragraphs[0].space_after = Pt(12)
    
    limitations = [
        '样本量仅11人/22眼，统计效力不足',
        '无对照组，无法准确计算控制率',
        '6个月随访偏短，需1~2年数据验证',
        '2例脱落，存在选择偏倚风险',
    ]
    
    for lim in limitations:
        lp = ltf.add_paragraph()
        lp.text = f'⚠  {lim}'
        lp.font.size = Pt(13)
        lp.font.color.rgb = COLORS['dark']
        lp.space_before = Pt(8)
    
    # 建议卡片
    rec_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.3), Inches(6), Inches(5.5))
    rec_card.fill.solid()
    rec_card.fill.fore_color.rgb = RGBColor(232, 245, 233)
    rec_card.line.color.rgb = COLORS['success']
    
    rtf = rec_card.text_frame
    rtf.word_wrap = True
    rtf.paragraphs[0].text = '后续建议'
    rtf.paragraphs[0].font.size = Pt(20)
    rtf.paragraphs[0].font.bold = True
    rtf.paragraphs[0].font.color.rgb = COLORS['success']
    rtf.paragraphs[0].space_after = Pt(12)
    
    recommendations = [
        '继续追踪至12个月及24个月',
        '增设对照组或引用历史对照数据',
        '对控制欠佳者进行个案分析',
        '加入角膜曲率数据',
        '增加脉络膜厚度数据（如条件允许）',
    ]
    
    for i, rec in enumerate(recommendations):
        rp = rtf.add_paragraph()
        rp.text = f'{i+1}.  {rec}'
        rp.font.size = Pt(13)
        rp.font.color.rgb = COLORS['dark']
        rp.space_before = Pt(8)
    
    add_page_number(slide, 10, 10)
    add_footer_text(slide)


# ==================== 主程序 ====================
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    print('开始生成PPT...')
    
    create_cover_slide(prs)
    print('  [1/10] 封面页')
    
    create_toc_slide(prs)
    print('  [2/10] 目录页')
    
    create_overview_slide(prs)
    print('  [3/10] 研究概况')
    
    create_al_trend_slide(prs)
    print('  [4/10] 眼轴增长趋势')
    
    create_ser_slide(prs)
    print('  [5/10] SER变化分析')
    
    create_subgroup_slide(prs)
    print('  [6/10] 亚组分析')
    
    create_control_slide(prs)
    print('  [7/10] 疗效控制率')
    
    create_comparison_slide(prs)
    print('  [8/10] 文献基准对比')
    
    create_satisfaction_slide(prs)
    print('  [9/10] 满意度评价')
    
    create_limitation_slide(prs)
    print('  [10/10] 局限性与建议')
    
    prs.save(OUTPUT_FILE)
    print(f'\nPPT已保存: {OUTPUT_FILE}')


if __name__ == '__main__':
    main()
